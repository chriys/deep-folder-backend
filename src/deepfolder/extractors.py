import asyncio
from io import BytesIO
from typing import Any

from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload


class PDFExtractor:
    @staticmethod
    async def extract_text(file_content: bytes) -> dict[int, str]:
        """Extract text from PDF by page. Returns dict of {page_num: text}."""
        try:
            import pypdf
        except ImportError:
            raise ImportError("pypdf is required for PDF extraction") from None

        def _extract() -> dict[int, str]:
            reader = pypdf.PdfReader(BytesIO(file_content))
            pages = {}
            for page_num, page in enumerate(reader.pages, 1):
                pages[page_num] = page.extract_text()
            return pages

        return await asyncio.to_thread(_extract)


class GoogleSlidesExtractor:
    @staticmethod
    async def extract_slides(
        file_id: str, credentials: Credentials
    ) -> dict[str, str]:
        """Extract text from Google Slides, keyed by objectId."""
        service = build("slides", "v1", credentials=credentials)

        def _extract() -> dict[str, str]:
            presentation = service.presentations().get(presentationId=file_id).execute()
            slides_dict: dict[str, str] = {}
            for slide in presentation.get("slides", []):
                object_id = slide.get("objectId", "")
                text_parts: list[str] = []
                for element in slide.get("pageElements", []):
                    if "shape" in element:
                        shape = element["shape"]
                        if "text" in shape:
                            for text_element in shape["text"].get("textElements", []):
                                if "textRun" in text_element:
                                    text_parts.append(text_element["textRun"]["content"])
                slides_dict[object_id] = "".join(text_parts).strip()
            return slides_dict

        return await asyncio.to_thread(_extract)


class GoogleSheetsExtractor:
    @staticmethod
    async def extract_sheets(
        file_id: str, credentials: Credentials
    ) -> list[dict[str, str]]:
        """Extract text from Google Sheets, one entry per sheet.
        Returns list of {name, gid, text, row_range} dicts.
        """
        service = build("sheets", "v4", credentials=credentials)

        def _extract() -> list[dict[str, str]]:
            spreadsheet = service.spreadsheets().get(spreadsheetId=file_id).execute()
            sheets_data: list[dict[str, str]] = []
            for sheet in spreadsheet.get("sheets", []):
                props = sheet.get("properties", {})
                sheet_title = props.get("title", "Sheet1")
                sheet_id = props.get("sheetId", 0)
                grid = props.get("gridProperties", {})
                row_count = grid.get("rowCount", 0)
                col_count = grid.get("columnCount", 0)

                if row_count == 0 or col_count == 0:
                    continue

                last_col = _column_letter(col_count)
                range_name = f"{sheet_title}!A1:{last_col}{row_count}"

                result = service.spreadsheets().values().get(
                    spreadsheetId=file_id, range=range_name
                ).execute()

                values = result.get("values", [])
                text_lines = ["\t".join(str(c) for c in row) for row in values]
                full_text = "\n".join(text_lines)

                sheets_data.append({
                    "name": sheet_title,
                    "gid": str(sheet_id),
                    "text": full_text,
                    "row_range": f"A1:{last_col}{row_count}",
                })
            return sheets_data

        return await asyncio.to_thread(_extract)


def _column_letter(n: int) -> str:
    """Convert 1-based column index to spreadsheet column letter (1=A, 26=Z, 27=AA)."""
    result = ""
    while n > 0:
        n -= 1
        result = chr(ord("A") + n % 26) + result
        n //= 26
    return result


class GoogleDocsExtractor:
    @staticmethod
    async def extract_text(
        file_id: str, credentials: Credentials
    ) -> str:
        """Extract text from Google Doc using export."""
        service = build("drive", "v3", credentials=credentials)

        def _extract() -> str:
            try:
                request = service.files().export_media(
                    fileId=file_id, mimeType="text/plain"
                )
                fh = BytesIO()
                downloader = MediaIoBaseDownload(fh, request)
                done = False
                while not done:
                    _, done = downloader.next_chunk()
                return fh.getvalue().decode("utf-8")
            except Exception as e:
                raise ValueError(f"Error extracting Google Doc {file_id}: {e}") from e

        return await asyncio.to_thread(_extract)

    @staticmethod
    async def extract_with_headings(
        file_id: str, credentials: Credentials
    ) -> tuple[str, list[dict[str, str]]]:
        """Extract text and heading structure from Google Doc.
        Returns (text, list of {text, anchor_id} dicts)."""
        service = build("docs", "v1", credentials=credentials)

        def _extract() -> tuple[str, list[dict[str, str]]]:
            try:
                doc = service.documents().get(documentId=file_id).execute()
                text = GoogleDocsExtractor._extract_text_from_document(doc)
                headings = GoogleDocsExtractor._extract_headings_from_document(doc)
                return text, headings
            except Exception as e:
                raise ValueError(f"Error extracting Google Doc {file_id}: {e}") from e

        return await asyncio.to_thread(_extract)

    @staticmethod
    def _extract_text_from_document(doc: dict[str, Any]) -> str:
        """Extract plain text from Google Doc structure."""
        text_parts = []
        for element in doc.get("body", {}).get("content", []):
            if "paragraph" in element:
                para = element["paragraph"]
                for run in para.get("elements", []):
                    if "textRun" in run:
                        text_parts.append(run["textRun"]["content"])
        return "".join(text_parts)

    @staticmethod
    def _extract_headings_from_document(doc: dict[str, Any]) -> list[dict[str, str]]:
        """Extract headings with anchor IDs from Google Doc structure."""
        headings: list[dict[str, str]] = []
        for element in doc.get("body", {}).get("content", []):
            if "paragraph" in element:
                para = element["paragraph"]
                style = para.get("paragraphStyle", {})
                heading_id = style.get("headingId")

                if heading_id:
                    text_parts = []
                    for run in para.get("elements", []):
                        if "textRun" in run:
                            text_parts.append(run["textRun"]["content"])
                    heading_text = "".join(text_parts).strip()
                    if heading_text:
                        headings.append({"text": heading_text, "anchor_id": heading_id})

        return headings


class DocxExtractor:
    @staticmethod
    async def extract_text(file_content: bytes) -> list[tuple[str, str]]:
        """Extract text from docx by heading sections.
        Returns list of (heading_text, section_content) tuples."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx is required for docx extraction") from None

        def _extract() -> list[tuple[str, str]]:
            doc = Document(BytesIO(file_content))
            sections: list[tuple[str, str]] = []
            current_heading = "Document"
            current_content: list[str] = []

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                style = para.style
                if style and style.name.startswith("Heading"):
                    if current_content:
                        sections.append((current_heading, "\n".join(current_content)))
                    current_heading = text
                    current_content = []
                else:
                    current_content.append(text)

            if current_content:
                sections.append((current_heading, "\n".join(current_content)))

            if not sections:
                full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                sections = [("Document", full_text)]

            return sections

        return await asyncio.to_thread(_extract)


class PptxExtractor:
    @staticmethod
    async def extract_text(file_content: bytes) -> dict[int, str]:
        """Extract text from pptx by slide. Returns dict of {slide_num: text}."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for pptx extraction") from None

        def _extract() -> dict[int, str]:
            prs = Presentation(BytesIO(file_content))
            slides: dict[int, str] = {}
            for slide_num, slide in enumerate(prs.slides, 1):
                texts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                slides[slide_num] = "\n".join(texts)
            return slides

        return await asyncio.to_thread(_extract)


class XlsxExtractor:
    @staticmethod
    async def extract_text(file_content: bytes) -> dict[str, str]:
        """Extract text from xlsx by sheet. Returns dict of {sheet_name: text}."""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required for xlsx extraction") from None

        def _extract() -> dict[str, str]:
            wb = openpyxl.load_workbook(BytesIO(file_content), data_only=True)
            sheets: dict[str, str] = {}
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows: list[str] = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    row_text = "\t".join(cells).strip()
                    if row_text:
                        rows.append(row_text)
                sheets[sheet_name] = "\n".join(rows)
            return sheets

        return await asyncio.to_thread(_extract)
