from io import BytesIO

import pytest

from deepfolder.extractors import DocxExtractor, PptxExtractor, XlsxExtractor


def _make_docx(heading_texts: list[str], body_texts: list[str]) -> bytes:
    """Create a minimal docx in memory and return bytes."""
    from docx import Document

    doc = Document()
    for heading, body in zip(heading_texts, body_texts):
        doc.add_heading(heading, level=1)
        doc.add_paragraph(body)
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx(slide_texts: list[str]) -> bytes:
    """Create a minimal pptx in memory and return bytes."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for text in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        tf = txBox.text_frame
        tf.text = text
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx(sheet_data: dict[str, list[list[str]]]) -> bytes:
    """Create a minimal xlsx in memory and return bytes."""
    import openpyxl

    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    for sheet_name, rows in sheet_data.items():
        ws = wb.create_sheet(title=sheet_name)
        for row in rows:
            ws.append(row)
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


class TestDocxExtractor:
    @pytest.mark.asyncio
    async def test_extract_with_headings(self) -> None:
        content = _make_docx(
            heading_texts=["Introduction", "Chapter 1"],
            body_texts=["Intro text here.", "Chapter content."],
        )
        sections = await DocxExtractor.extract_text(content)

        assert len(sections) >= 2
        assert sections[0][0] == "Introduction"
        assert "Intro text here." in sections[0][1]
        assert sections[1][0] == "Chapter 1"
        assert "Chapter content." in sections[1][1]

    @pytest.mark.asyncio
    async def test_extract_no_headings(self) -> None:
        from docx import Document

        doc = Document()
        doc.add_paragraph("Just some body text.")
        doc.add_paragraph("More text without headings.")
        content = BytesIO()
        doc.save(content)
        sections = await DocxExtractor.extract_text(content.getvalue())

        assert len(sections) == 1
        assert sections[0][0] == "Document"
        assert "Just some body text." in sections[0][1]

    @pytest.mark.asyncio
    async def test_extract_empty_document(self) -> None:
        from docx import Document

        doc = Document()
        content = BytesIO()
        doc.save(content)
        sections = await DocxExtractor.extract_text(content.getvalue())

        assert len(sections) == 1
        assert sections[0][1] == ""


class TestPptxExtractor:
    @pytest.mark.asyncio
    async def test_extract_slides(self) -> None:
        content = _make_pptx(["Slide 1 content", "Slide 2 content"])
        slides = await PptxExtractor.extract_text(content)

        assert len(slides) == 2
        assert slides[1] == "Slide 1 content"
        assert slides[2] == "Slide 2 content"

    @pytest.mark.asyncio
    async def test_extract_empty_slides(self) -> None:
        from pptx import Presentation

        prs = Presentation()
        buf = BytesIO()
        prs.save(buf)
        slides = await PptxExtractor.extract_text(buf.getvalue())

        assert len(slides) == 0


class TestXlsxExtractor:
    @pytest.mark.asyncio
    async def test_extract_sheets(self) -> None:
        content = _make_xlsx({"Sheet1": [["A1", "B1"], ["A2", "B2"]]})
        sheets = await XlsxExtractor.extract_text(content)

        assert len(sheets) == 1
        assert "Sheet1" in sheets
        assert "A1\tB1" in sheets["Sheet1"]
        assert "A2\tB2" in sheets["Sheet1"]

    @pytest.mark.asyncio
    async def test_extract_multiple_sheets(self) -> None:
        content = _make_xlsx({
            "Sheet1": [["Data1"]],
            "Sheet2": [["Data2"]],
        })
        sheets = await XlsxExtractor.extract_text(content)

        assert len(sheets) == 2
        assert "Sheet1" in sheets
        assert "Sheet2" in sheets

    @pytest.mark.asyncio
    async def test_extract_empty_sheet(self) -> None:
        import openpyxl

        wb = openpyxl.Workbook()
        buf = BytesIO()
        wb.save(buf)
        sheets = await XlsxExtractor.extract_text(buf.getvalue())

        assert len(sheets) >= 1
        assert sheets.get("Sheet") is not None
        assert sheets["Sheet"] == ""
